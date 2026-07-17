"""Renderer contract tests. HTML tests are stdlib-only so they run on any host;
PPTX tests skip when python-pptx is absent (it lives in the contained computer)."""
import importlib.util
import os
import re
import shutil
import tempfile
import unittest

HERE = os.path.dirname(os.path.abspath(__file__))
_spec = importlib.util.spec_from_file_location("deck_render", os.path.join(HERE, "deck_render.py"))
deck_render = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(deck_render)

_qa_spec = importlib.util.spec_from_file_location("deck_qa", os.path.join(HERE, "deck_qa.py"))
deck_qa = importlib.util.module_from_spec(_qa_spec)
_qa_spec.loader.exec_module(deck_qa)


def _find_chromium():
    """Locate a Chromium/Chrome binary so the rendered-QA tests can run where one
    exists (dev macs, the contained computer) and skip cleanly where it does not
    (headless CI without a browser). CHROMIUM env wins, then PATH, then the
    standard macOS app bundles."""
    env = os.environ.get("CHROMIUM")
    if env and os.path.exists(env):
        return env
    for name in ("chromium", "chromium-browser", "google-chrome", "google-chrome-stable", "chrome"):
        found = shutil.which(name)
        if found:
            return found
    for bundle in (
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ):
        if os.path.exists(bundle):
            return bundle
    return None


CHROMIUM = _find_chromium()

ALL_LAYOUTS_DECK = {
    "title": "T",
    "slides": [
        {"layout": "cover", "title": "T", "subtitle": "S"},
        {"layout": "timeline", "title": "Roadmap", "items": [
            {"label": "Q3", "title": "Ship", "text": "Self-serve"},
            {"label": "Q4", "title": "Scale", "text": "EU launch"},
        ]},
        {"layout": "comparison", "title": "Risks", "headers": ["Risk", "Impact"],
         "rows": [["Fuel", "High"], ["Churn", "Low"]]},
        {"layout": "team_grid", "title": "Team", "members": [
            {"name": "Elena Ricci", "role": "CEO"},
            {"name": "Marco Chen", "role": "CTO"},
        ]},
        {"layout": "closing", "title": "Next", "bullets": ["a"]},
    ],
}


class RenderHtmlLayouts(unittest.TestCase):
    def test_new_layouts_render_and_css_formats(self):
        # Also guards the _HTML_CSS .format() contract: an unescaped { in new CSS
        # raises KeyError/IndexError here.
        html = deck_render.render_html(ALL_LAYOUTS_DECK, HERE)
        self.assertIn('class="tl-item"', html)
        self.assertIn("Q3", html)
        self.assertIn('<table class="cmp">', html)
        self.assertIn("<th>Risk</th>", html)
        self.assertIn('class="member"', html)
        self.assertIn(">ER<", html)  # initials avatar for Elena Ricci

    def test_initials(self):
        self.assertEqual(deck_render._initials("Elena Ricci"), "ER")
        self.assertEqual(deck_render._initials("Cher"), "C")
        self.assertEqual(deck_render._initials(""), "")

    def test_theme_fonts_are_css_variables(self):
        # Pins the parametric-font contract (--head/--body defined AND consumed):
        # the brand-kit live recolor injects var overrides, so a regression to
        # baked-in font literals would silently break recolor on deck previews.
        html = deck_render.render_html(
            {"title": "T", "theme": {"heading_font": "Georgia"},
             "slides": [{"layout": "cover", "title": "T"}]}, HERE)
        self.assertIn("--head:'Georgia'", html)
        self.assertIn("var(--head)", html)
        self.assertIn("var(--body)", html)


class EditorialCover(unittest.TestCase):
    def test_cover_renders_eyebrow_and_hero_art(self):
        html = deck_render.render_html(
            {"title": "T", "theme": {"name": "editorial_bold"},
             "slides": [{"layout": "cover", "title": "Kite", "subtitle": "S",
                         "eyebrow": "EyebrowProbe", "hero_art": "rings"}]}, HERE)
        self.assertIn("EyebrowProbe", html)
        self.assertIn("hero-art", html)          # procedural svg wrapper class
        self.assertIn("--surface:#0f3d3e", html)  # theme surface reaches :root
        # S1a final-review Fix 5: aria-hidden must be a valid boolean attribute,
        # not the bare/invalid `aria-hidden` HTML shorthand.
        self.assertIn('aria-hidden="true"', html)
        self.assertNotIn("aria-hidden>", html)
        self.assertNotIn("aria-hidden ", html)

    def test_surface_ink_reach_root_for_all_themes(self):
        html = deck_render.render_html(
            {"title": "T", "theme": {"name": "editorial_noir"},
             "slides": [{"layout": "cover", "title": "X"}]}, HERE)
        self.assertIn("--surface:#0b0b0d", html)
        self.assertIn("--ink:#f4f1ea", html)

    def test_hero_art_grid_ids_are_unique_across_slides(self):
        # Review gotcha: _hero_art("grid") used a FIXED pattern id — two grid
        # slides (e.g. cover + section) in one deck produced duplicate DOM
        # ids. Each call must mint its own id.
        html = deck_render.render_html(
            {"title": "T", "slides": [
                {"layout": "cover", "title": "A", "hero_art": "grid"},
                {"layout": "section", "title": "B", "hero_art": "grid"},
            ]}, HERE)
        pattern_ids = re.findall(r'<pattern id="(g[^"]*)"', html)
        self.assertEqual(len(pattern_ids), 2)
        self.assertEqual(len(pattern_ids), len(set(pattern_ids)))
        for pid in pattern_ids:
            self.assertIn(f'url(#{pid})', html)


@unittest.skipUnless(
    importlib.util.find_spec("pptx"), "python-pptx not installed on this host"
)
class RenderPptxLayouts(unittest.TestCase):
    @staticmethod
    def _slide_texts(slide):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        return " ".join(texts)

    def test_new_layouts_produce_slides(self):
        import tempfile
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as tmp:
            out = os.path.join(tmp, "deck.pptx")
            stats = deck_render.render_pptx(ALL_LAYOUTS_DECK, tmp, out)
            self.assertIsNotNone(stats)
            prs = Presentation(out)
            self.assertEqual(len(prs.slides), len(ALL_LAYOUTS_DECK["slides"]))

            # Per-layout content assertions: catch a dead/mistyped elif branch that
            # would otherwise silently fall through to the bullets fallback and
            # still pass a slide-count-only check. Order follows ALL_LAYOUTS_DECK:
            # cover, timeline, comparison, team_grid, closing.
            # GOTCHA: the probe text must NOT be a substring of the slide TITLE —
            # the title renders via the shared path for EVERY layout, so a
            # title-substring probe (e.g. "Risk" ⊂ "Risks") passes vacuously.
            slides = list(prs.slides)
            self.assertIn("Q3", self._slide_texts(slides[1]))       # timeline label rendered
            self.assertIn("Fuel", self._slide_texts(slides[2]))     # comparison body cell (not in title)
            self.assertTrue(any(getattr(shape, "has_table", False)
                                for shape in slides[2].shapes))     # a real table shape exists
            self.assertIn("ER", self._slide_texts(slides[3]))       # team_grid initials avatar

    @staticmethod
    def _cover_fill_and_text_hex(prs):
        """Locate the cover's full-bleed background shape (fill hex) and the
        title run's font colour (text hex), independent of shape order."""
        slide = prs.slides[0]
        bg = [s for s in slide.shapes
              if getattr(s, "width", None) == prs.slide_width
              and getattr(s, "height", None) == prs.slide_height]
        title = next(
            s for s in slide.shapes
            if getattr(s, "has_text_frame", False) and s.text_frame.text.startswith("Kite")
        )
        run = title.text_frame.paragraphs[0].runs[0]
        return str(bg[0].fill.fore_color.rgb), str(run.font.color.rgb)

    def test_editorial_cover_uses_secondary_fill_and_ink_text(self):
        # S1a final-review Fix 2: fill=primary + hardcoded white text made
        # editorial_noir's gold-on-gold cover illegible and diverged from the
        # HTML/PDF preview (near-black surface, cream text). The pptx cover
        # must now match render_html's --surface/--ink for editorial themes.
        import tempfile
        from pptx import Presentation
        deck = {
            "title": "T",
            "theme": {"name": "editorial_noir", "primary": "#c9a54e", "secondary": "#1a1a1e",
                      "accent": "#c9a54e", "heading_font": "Georgia", "body_font": "Inter"},
            "slides": [{"layout": "cover", "title": "Kite", "subtitle": "Sub"}],
        }
        with tempfile.TemporaryDirectory() as tmp:
            out = os.path.join(tmp, "deck.pptx")
            deck_render.render_pptx(deck, tmp, out)
            fill_hex, text_hex = self._cover_fill_and_text_hex(Presentation(out))
        self.assertEqual(fill_hex, "1A1A1E")  # theme secondary (dark)
        self.assertEqual(text_hex, "F4F1EA")  # theme ink (cream) — NOT hardcoded white

    def test_editorial_light_theme_cover_keeps_dark_ink_text(self):
        # The other half of the same fix: a LIGHT editorial theme (surface AND
        # secondary are light) must NOT get hardcoded white text — that would
        # be as illegible as the noir bug this fix closes, just inverted.
        import tempfile
        from pptx import Presentation
        deck = {
            "title": "T",
            "theme": {"name": "editorial_ivory", "primary": "#1f4d3f", "secondary": "#e9e3d6",
                      "accent": "#1f4d3f", "heading_font": "Georgia", "body_font": "Inter"},
            "slides": [{"layout": "cover", "title": "Kite", "subtitle": "Sub"}],
        }
        with tempfile.TemporaryDirectory() as tmp:
            out = os.path.join(tmp, "deck.pptx")
            deck_render.render_pptx(deck, tmp, out)
            fill_hex, text_hex = self._cover_fill_and_text_hex(Presentation(out))
        self.assertEqual(fill_hex, "E9E3D6")  # theme secondary (light)
        self.assertEqual(text_hex, "1C1B18")  # theme ink (dark) — NOT hardcoded white

    def test_legacy_theme_cover_unchanged_fill_primary_text_white(self):
        # Legacy (non-editorial) themes must NOT be touched by Fix 2: their
        # `ink` is tuned for a white `surface`, not for their own (also dark)
        # `secondary` — reusing it here would print dark-on-dark.
        import tempfile
        from pptx import Presentation
        deck = {
            "title": "T",
            "theme": {"name": "high_contrast", "primary": "#111827", "secondary": "#000000",
                      "accent": "#f59e0b", "heading_font": "Inter", "body_font": "Inter"},
            "slides": [{"layout": "cover", "title": "Kite", "subtitle": "Sub"}],
        }
        with tempfile.TemporaryDirectory() as tmp:
            out = os.path.join(tmp, "deck.pptx")
            deck_render.render_pptx(deck, tmp, out)
            fill_hex, text_hex = self._cover_fill_and_text_hex(Presentation(out))
        self.assertEqual(fill_hex, "111827")  # theme primary, unchanged
        self.assertEqual(text_hex, "FFFFFF")  # hardcoded white, unchanged


class DeckQaOverflow(unittest.TestCase):
    """Guards deck_qa's slide_overflow check against the hero_art false positive:
    a decorative accent (.hero-art: right:-4vw;width:44vw) intentionally bleeds
    past the slide edge and is clipped by .slide{overflow:hidden}, but
    scrollWidth still measures its unclipped box. The check must ignore such
    intentionally-clipped decorative layers while still catching real overflow."""

    def _run_qa_on_html(self, html):
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "preview.html")
            with open(path, "w", encoding="utf-8") as handle:
                handle.write(html)
            result = deck_qa.run_qa(path, CHROMIUM, mode="deck")
        return [issue.get("code") for issue in result.get("issues", [])], result

    def test_qa_js_excludes_decorative_layers_from_overflow(self):
        # Browser-free canary: runs in CI even without Chromium. If a refactor
        # drops the decorative-layer exclusion, scrollWidth would false-flag
        # every hero_art cover/section again — catch it without a browser.
        self.assertIn("hero-art", deck_qa.QA_JS)
        self.assertIn("pointerEvents", deck_qa.QA_JS)

    @unittest.skipUnless(CHROMIUM, "no chromium/chrome binary found")
    def test_hero_art_bleed_not_flagged_as_overflow(self):
        # End-to-end: the real renderer output for a hero_art cover must pass the
        # overflow check. This is the exact scenario that regressed on every deck
        # pack once hero_art got real values.
        html = deck_render.render_html(
            {"title": "T", "theme": {"name": "editorial_bold"},
             "slides": [{"layout": "cover", "title": "Kite", "subtitle": "S",
                         "eyebrow": "E", "hero_art": "rings"}]}, HERE)
        codes, result = self._run_qa_on_html(html)
        self.assertNotIn(
            "slide_overflow", codes,
            f"hero_art decorative bleed must not be flagged as overflow: {result.get('issues')}")

    @unittest.skipUnless(CHROMIUM, "no chromium/chrome binary found")
    def test_real_content_overflow_still_flagged(self):
        # Guards against an over-broad fix: a genuinely too-wide, non-decorative
        # element (default pointer-events, no hero-art class) must STILL trip the
        # check — the exclusion is surgical to decorative layers only.
        html = (
            "<!doctype html><html><head><meta charset='utf-8'><style>"
            ".slide{width:1280px;height:720px;position:relative;overflow:hidden;box-sizing:border-box}"
            ".wide{width:1600px;height:40px;background:#333}"
            "</style></head><body>"
            "<section class='slide'><div class='wide'></div></section></body></html>"
        )
        codes, result = self._run_qa_on_html(html)
        self.assertIn(
            "slide_overflow", codes,
            f"genuine content overflow must still be flagged: {result.get('issues')}")


if __name__ == "__main__":
    unittest.main()
