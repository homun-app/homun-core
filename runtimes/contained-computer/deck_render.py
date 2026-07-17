#!/usr/bin/env python3
"""Homun deck renderer — ONE structured deck model → TWO outputs.

The agent produces only the CONTENT as a JSON deck (theme + slides with a fixed set
of layout types); this renderer turns it into:

  * <prefix>.html  — a self-contained, on-brand, image-led deck for VIEWING / PDF
                     (logo + images inlined as data URLs, brand colours/fonts);
  * <prefix>.pptx  — a NATIVE, EDITABLE PowerPoint (real text boxes, logo, images)
                     built from scratch with python-pptx — NOT an HTML screenshot.

This separation (single source of truth → dual render) is what makes the deck both
beautiful AND editable in PowerPoint/Google Slides, instead of the lossy HTML→PPTX
trap. Layout types are deliberately CONSTRAINED so both renderers stay faithful.

Usage:
    python deck_render.py deck.json [--prefix deck] [--no-pptx]
    python deck_render.py deck.json [--prefix deck] [--template-pptx source.pptx]

deck.json schema (all fields optional unless noted):
{
  "title": "Deck title",            "subtitle": "...", "date": "...",
  "theme": {                         # falls back to brand-kit-ish defaults
    "name": "editorial_bold",         # optional: resolves surface/ink/muted/hairline/
                                       # on_brand from design_tokens.THEMES; explicit
                                       # fields below still override the named theme
    "primary": "#2b6cb0", "secondary": "#1a202c", "accent": "#ed8936",
    "heading_font": "Inter", "body_font": "Inter",
    "logo": "logo.png"               # path (rel to deck.json) or data: URL
  },
  "slides": [                        # REQUIRED
    {"layout": "cover",   "title": "...", "subtitle": "...",
     "eyebrow": "SEED ROUND",         # optional small-caps kicker above the title
     "hero_art": "rings"},            # optional procedural SVG accent: rings|grid|gradient|none
    {"layout": "bullets", "title": "...", "bullets": ["a","b"], "notes": "..."},
    {"layout": "image_right", "title": "...", "bullets": [...], "image": "s2.png"},
    {"layout": "kpi",     "title": "...", "kpi": "42%", "kpi_label": "growth"},
    {"layout": "two_column", "title": "...", "columns": [{"title":"L","bullets":[]},
                                                          {"title":"R","bullets":[]}]},
    {"layout": "quote",   "quote": "...", "author": "..."},
    {"layout": "section", "title": "...", "eyebrow": "...", "hero_art": "grid"},
    {"layout": "timeline", "title": "...", "items": [{"label":"Q1","title":"...","text":"..."}]},
    {"layout": "comparison", "title": "...", "headers": ["A","B"], "rows": [["x","y"]]},
    {"layout": "team_grid", "title": "...", "members": [{"name":"...","role":"..."}]},
    {"layout": "closing", "title": "Next steps", "bullets": [...]}
  ]
}
"""

import argparse
import base64
import html
import json
import mimetypes
import os
import re
import sys

from design_tokens import theme_values

DEFAULT_THEME = {
    "primary": "#2b6cb0",
    "secondary": "#1a202c",
    "accent": "#ed8936",
    "heading_font": "Inter",
    "body_font": "Inter",
    "logo": "",
}


# --------------------------------------------------------------------------- utils
def hexrgb(value, fallback=(43, 108, 176)):
    """'#rrggbb' → (r,g,b) ints, tolerant of missing '#'/bad input."""
    try:
        v = value.lstrip("#")
        if len(v) == 3:
            v = "".join(c * 2 for c in v)
        return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except Exception:
        return fallback


def data_url(path, base_dir):
    """Inline a file as a data: URL. Accepts an already-data: string. '' on miss."""
    if not path:
        return ""
    if path.startswith("data:"):
        return path
    full = path if os.path.isabs(path) else os.path.join(base_dir, path)
    if not os.path.isfile(full):
        return ""
    mime = mimetypes.guess_type(full)[0] or "image/png"
    with open(full, "rb") as fh:
        b64 = base64.b64encode(fh.read()).decode("ascii")
    return f"data:{mime};base64,{b64}"


def resolve_path(path, base_dir):
    if not path or path.startswith("data:"):
        return ""
    full = path if os.path.isabs(path) else os.path.join(base_dir, path)
    return full if os.path.isfile(full) else ""


# ----------------------------------------------------------------------- HTML side
def html_escape(s):
    return html.escape(str(s or ""))


def render_html(deck, base_dir):
    raw_theme = deck.get("theme") or {}
    name = raw_theme.get("name")
    # Mirrors doc_render.render_html: a NAMED theme resolves through design_tokens
    # (bringing surface/ink/muted/hairline/on_brand along), explicit fields still
    # win. With no name, keep the legacy behaviour — only DEFAULT_THEME + explicit
    # overrides — so old brand-kit-only decks don't silently inherit
    # clean_corporate's tokens (theme_values() defaults nameless lookups there).
    if name:
        theme = {**DEFAULT_THEME, **theme_values(name, raw_theme)}
    else:
        theme = {**DEFAULT_THEME, **{k: v for k, v in raw_theme.items() if v}}
    logo = data_url(theme.get("logo", ""), base_dir)
    slides_html = []
    for s in deck.get("slides", []):
        slides_html.append(_html_slide(s, base_dir, logo))
    css = _HTML_CSS.format(
        primary=theme["primary"],
        secondary=theme["secondary"],
        accent=theme["accent"],
        heading=theme["heading_font"],
        body=theme["body_font"],
        surface=theme.get("surface", "#ffffff"),
        ink=theme.get("ink", "#16202b"),
        muted=theme.get("muted", "#5a6675"),
        hairline=theme.get("hairline", "#e4e9ef"),
        on_brand=theme.get("on_brand", "#ffffff"),
    )
    title = html_escape(deck.get("title", "Presentation"))
    return _HTML_SHELL.format(title=title, css=css, body="\n".join(slides_html))


def _bullets_html(items):
    if not items:
        return ""
    lis = "".join(f"<li>{html_escape(b)}</li>" for b in items)
    return f'<ul class="body">{lis}</ul>'


def _logo_html(logo):
    return f'<img class="logo" src="{logo}">' if logo else ""


def _eyebrow(text):
    return f'<div class="eyebrow">{html_escape(text)}</div>' if text else ""


def _hero_art(kind):
    # Procedural editorial art — inline SVG, zero external images (license-clean,
    # local). Uses currentColor so it inherits the accent set on the cover.
    if kind == "rings":
        return ('<svg class="hero-art" viewBox="0 0 400 400" aria-hidden><g fill="none" '
                'stroke="currentColor" stroke-width="1.5" opacity=".5">'
                + "".join(f'<circle cx="300" cy="90" r="{r}"/>' for r in (40, 80, 120, 170))
                + "</g></svg>")
    if kind == "grid":
        return ('<svg class="hero-art" viewBox="0 0 400 400" aria-hidden>'
                '<defs><pattern id="g" width="26" height="26" patternUnits="userSpaceOnUse">'
                '<path d="M26 0H0V26" fill="none" stroke="currentColor" stroke-width="1" '
                'opacity=".35"/></pattern></defs><rect width="400" height="400" fill="url(#g)"/></svg>')
    if kind == "gradient":
        return '<div class="hero-art hero-grad" aria-hidden></div>'
    return ""


def _initials(name):
    """First letter of up to 2 words → avatar glyph when no headshot is provided."""
    parts = [p for p in str(name or "").split() if p]
    return "".join(p[0].upper() for p in parts[:2])


def _html_slide(s, base_dir, logo):
    layout = s.get("layout", "bullets")
    title = html_escape(s.get("title", ""))
    img = data_url(s.get("image", ""), base_dir)
    if layout == "cover":
        return (
            f'<section class="slide cover">{_logo_html(logo)}'
            f'{_hero_art(s.get("hero_art", ""))}'
            f'{_eyebrow(s.get("eyebrow", ""))}'
            f"<h1>{title}</h1>"
            f'<div class="sub">{html_escape(s.get("subtitle",""))}</div>'
            f'<div class="rule"></div></section>'
        )
    if layout == "section":
        return (
            f'<section class="slide section">{_logo_html(logo)}'
            f'{_hero_art(s.get("hero_art", ""))}'
            f'{_eyebrow(s.get("eyebrow", ""))}'
            f"<h1>{title}</h1><div class=\"accent-bar\"></div></section>"
        )
    if layout == "kpi":
        return (
            f'<section class="slide kpi-slide">{_logo_html(logo)}'
            f"<h2>{title}</h2>"
            f'<div class="kpi">{html_escape(s.get("kpi",""))}</div>'
            f'<div class="sub">{html_escape(s.get("kpi_label",""))}</div>'
            f'<div class="accent-bar"></div></section>'
        )
    if layout == "quote":
        return (
            f'<section class="slide quote-slide">{_logo_html(logo)}'
            f'<blockquote>“{html_escape(s.get("quote",""))}”</blockquote>'
            f'<div class="sub">— {html_escape(s.get("author",""))}</div>'
            f'<div class="accent-bar"></div></section>'
        )
    if layout in ("image_left", "image_right"):
        img_tag = f'<img class="led" src="{img}">' if img else '<div class="led ph"></div>'
        text = f"<div><h2>{title}</h2>{_bullets_html(s.get('bullets'))}" + (
            f'<p class="body">{html_escape(s.get("body",""))}</p>' if s.get("body") else ""
        ) + "</div>"
        order = (img_tag + text) if layout == "image_left" else (text + img_tag)
        return (
            f'<section class="slide img-led">{_logo_html(logo)}{order}'
            f'<div class="accent-bar"></div></section>'
        )
    if layout == "two_column":
        cols = s.get("columns", [])[:2]
        cells = "".join(
            f'<div class="col"><h3>{html_escape(c.get("title",""))}</h3>'
            f"{_bullets_html(c.get('bullets'))}</div>"
            for c in cols
        )
        return (
            f'<section class="slide two-col">{_logo_html(logo)}'
            f"<h2>{title}</h2><div class=\"cols\">{cells}</div>"
            f'<div class="accent-bar"></div></section>'
        )
    if layout == "timeline":
        items = s.get("items", [])[:6]
        rows = "".join(
            f'<div class="tl-item"><div class="tl-label">{html_escape(i.get("label", ""))}</div>'
            f'<div class="tl-dot"></div>'
            f'<div class="tl-text"><strong>{html_escape(i.get("title", ""))}</strong>'
            f'<span>{html_escape(i.get("text", ""))}</span></div></div>'
            for i in items
        )
        return (
            f'<section class="slide timeline">{_logo_html(logo)}'
            f'<h2>{title}</h2><div class="tl">{rows}</div>'
            f'<div class="accent-bar"></div></section>'
        )
    if layout == "comparison":
        headers = s.get("headers", [])[:4]
        rows = s.get("rows", [])[:8]
        table = ""
        # Mirror the PPTX guard: a headerless/rowless table is worse than no table
        # (an empty <thead>/<tbody> still renders borders and wastes the slide).
        if headers and rows:
            head = "".join(f"<th>{html_escape(h)}</th>" for h in headers)
            body_rows = "".join(
                "<tr>" + "".join(f"<td>{html_escape(c)}</td>" for c in row[: len(headers)]) + "</tr>"
                for row in rows
            )
            table = f'<table class="cmp"><thead><tr>{head}</tr></thead><tbody>{body_rows}</tbody></table>'
        return (
            f'<section class="slide comparison">{_logo_html(logo)}'
            f"<h2>{title}</h2>{table}"
            f'<div class="accent-bar"></div></section>'
        )
    if layout == "team_grid":
        members = s.get("members", [])[:8]
        cells = "".join(
            f'<div class="member"><div class="avatar">{html_escape(_initials(m.get("name", "")))}</div>'
            f'<strong>{html_escape(m.get("name", ""))}</strong>'
            f'<span>{html_escape(m.get("role", ""))}</span></div>'
            for m in members
        )
        return (
            f'<section class="slide team">{_logo_html(logo)}'
            f'<h2>{title}</h2><div class="team-grid">{cells}</div>'
            f'<div class="accent-bar"></div></section>'
        )
    # default: bullets (+ optional image, + optional body)
    body = f'<p class="body">{html_escape(s.get("body",""))}</p>' if s.get("body") else ""
    img_block = f'<img class="inline-img" src="{img}">' if img else ""
    return (
        f'<section class="slide bullets">{_logo_html(logo)}'
        f"<h2>{title}</h2>{_bullets_html(s.get('bullets'))}{body}{img_block}"
        f'<div class="accent-bar"></div></section>'
    )


_HTML_CSS = """
:root{{--brand:{primary};--brand2:{secondary};--accent:{accent};
  --head:'{heading}';--body:'{body}';
  --ink:{ink};--muted:{muted};--surface:{surface};--hairline:{hairline};--on-brand:{on_brand};}}
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:var(--head),-apple-system,'Segoe UI',Roboto,Helvetica,Arial,sans-serif;
  color:var(--ink);-webkit-font-smoothing:antialiased}}
.slide{{width:100%;min-height:100vh;padding:9vh 9vw;display:flex;flex-direction:column;
  justify-content:center;position:relative;page-break-after:always;overflow:hidden;
  background:var(--surface)}}
.slide h1{{font-size:4rem;line-height:1.04;font-weight:800;letter-spacing:-.02em}}
.slide h2{{font-size:2.7rem;line-height:1.1;font-weight:800;letter-spacing:-.01em;
  color:var(--ink);margin-bottom:.7em;padding-bottom:.32em;
  border-bottom:4px solid var(--accent);display:inline-block}}
.slide h3{{font-size:1.4rem;color:var(--brand);font-weight:700;margin-bottom:.4em}}
.slide h1,.slide h2,.slide h3,.slide li,.slide p,.slide .sub,.slide .col,
.slide blockquote,.slide .kpi{{overflow-wrap:anywhere;hyphens:auto}}
.body{{font-family:var(--body),-apple-system,sans-serif}}
.slide ul{{margin-top:.4rem}}
.slide li{{font-size:1.6rem;line-height:1.4;color:var(--muted);margin:.75rem 0;
  list-style:none;padding-left:1.9rem;position:relative}}
.slide li::before{{content:"";position:absolute;left:0;top:.55em;width:.72rem;height:.72rem;
  border-radius:2px;background:var(--accent)}}
/* left accent rail on every content slide */
.slide:not(.cover):not(.section)::before{{content:"";position:absolute;left:0;top:0;bottom:0;
  width:12px;background:var(--brand)}}
.kpi{{font-size:7rem;font-weight:800;color:var(--brand);line-height:1;letter-spacing:-.03em}}
.accent-bar{{position:absolute;left:0;bottom:0;height:8px;width:100%;
  background:linear-gradient(90deg,var(--brand),var(--accent))}}
.logo{{position:absolute;top:6vh;right:9vw;max-height:46px}}
.cover .sub{{font-size:1.5rem;opacity:.92;margin-top:1.3rem;font-weight:400;position:relative}}
.cover .rule{{width:96px;height:6px;background:var(--accent);margin-top:2rem;position:relative}}
.img-led{{display:grid;grid-template-columns:1fr 1fr;gap:5vw;align-items:center}}
.img-led .led{{width:100%;border-radius:16px;object-fit:cover;max-height:64vh;
  box-shadow:0 14px 44px rgba(0,0,0,.16)}}
.img-led .ph{{background:var(--hairline);min-height:42vh;border-radius:16px}}
.inline-img{{margin-top:1.4rem;max-height:44vh;border-radius:14px;object-fit:cover;
  box-shadow:0 10px 30px rgba(0,0,0,.12)}}
.two-col .cols{{display:grid;grid-template-columns:1fr 1fr;gap:5vw;margin-top:1.2rem}}
.two-col .col{{border-top:4px solid var(--accent);padding-top:1rem}}
.kpi-slide .sub{{font-size:1.6rem;color:var(--muted);margin-top:.6rem}}
.quote-slide blockquote{{font-size:2.8rem;font-weight:700;color:var(--ink);max-width:86%;
  line-height:1.25}}
.quote-slide blockquote::first-letter{{color:var(--accent)}}
.quote-slide .sub{{font-size:1.4rem;color:var(--brand);margin-top:1.4rem;font-weight:600}}
@media print{{.slide{{min-height:auto;height:100vh}}}}
.tl{{display:flex;flex-direction:column;gap:1.15rem;margin-top:1.3rem}}
.tl-item{{display:grid;grid-template-columns:92px 18px 1fr;align-items:start;gap:1.1rem}}
.tl-label{{font-weight:800;color:var(--brand);font-size:1.15rem;text-align:right}}
.tl-dot{{width:14px;height:14px;border-radius:50%;background:var(--accent);margin-top:.28rem;position:relative}}
.tl-item:not(:last-child) .tl-dot::after{{content:"";position:absolute;left:6px;top:17px;width:2px;height:2.6rem;background:var(--hairline)}}
.tl-text strong{{font-size:1.25rem}}
.tl-text span{{display:block;color:var(--muted);font-size:1.05rem;margin-top:.18rem}}
table.cmp{{width:100%;border-collapse:collapse;margin-top:1.3rem;font-size:1.15rem}}
table.cmp th{{text-align:left;background:var(--brand);color:var(--on-brand);padding:.7rem .95rem;font-weight:700}}
table.cmp td{{padding:.68rem .95rem;color:var(--muted);border-bottom:1px solid var(--hairline)}}
table.cmp tr:nth-child(even) td{{background:color-mix(in srgb,var(--hairline) 40%,transparent)}}
.team-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(200px,1fr));gap:1.7rem;margin-top:1.5rem}}
.member{{display:flex;flex-direction:column;align-items:flex-start;gap:.4rem}}
.member .avatar{{width:64px;height:64px;border-radius:50%;background:var(--brand);color:var(--on-brand);display:flex;align-items:center;justify-content:center;font-weight:800;font-size:1.3rem}}
.member strong{{font-size:1.2rem}}
.member span{{color:var(--muted);font-size:1rem}}
.eyebrow{{text-transform:uppercase;letter-spacing:.28em;font-size:.95rem;font-weight:700;
  color:var(--accent);margin-bottom:1.1rem;position:relative}}
.cover,.section{{background:var(--surface);color:var(--ink)}}
.cover h1,.section h1{{color:var(--ink);font-size:5rem;max-width:88%;position:relative}}
.cover .sub{{color:var(--muted)}}
.hero-art{{position:absolute;right:-4vw;top:-4vw;width:44vw;height:44vw;color:var(--accent);
  opacity:.9;pointer-events:none}}
.hero-art.hero-grad{{background:radial-gradient(120% 120% at 80% 0%,
  color-mix(in srgb,var(--accent) 40%,transparent),transparent 60%)}}
"""

_HTML_SHELL = """<!doctype html><html lang="en"><head><meta charset="utf-8">
<title>{title}</title><style>{css}</style></head><body>
{body}
</body></html>"""


# ----------------------------------------------------------------------- PPTX side
def render_template_pptx(deck, base_dir, template_path, out_path):
    """Edit a real PPTX template in place, preserving its masters/media/layouts.

    This is intentionally conservative. Imported templates are visually stronger
    than Homun's synthetic layouts, so the first production pass changes only
    text content and trims text aggressively enough to avoid the worst overflow.
    """
    try:
        from pptx import Presentation
    except Exception as exc:
        sys.stderr.write(f"deck_render: python-pptx unavailable ({exc}); skipping template .pptx\n")
        return False

    source = template_path if os.path.isabs(template_path) else os.path.join(base_dir, template_path)
    if not os.path.isfile(source):
        sys.stderr.write(f"deck_render: template pptx not found: {template_path}\n")
        return False

    prs = Presentation(source)
    slides_list = deck.get("slides", [])
    if not slides_list:
        return False

    def trim(text, max_chars):
        text = " ".join(str(text or "").split())
        if len(text) <= max_chars:
            return text
        return text[: max(1, max_chars - 1)].rstrip() + "…"

    def normalized(text):
        return " ".join(str(text or "").split()).lower()

    def shape_text(shape):
        return " ".join(str(getattr(shape, "text", "") or "").split())

    def text_shapes(slide):
        shapes = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape_text(shape):
                shapes.append(shape)
        return shapes

    def set_shape_text(shape, text):
        text = str(text or "")
        if not text.strip():
            return
        # Prefer run-level edits when possible to preserve font styling. Fall
        # back to shape.text for grouped/split text where the source text spans
        # multiple runs and cannot be replaced safely.
        try:
            paragraphs = shape.text_frame.paragraphs
            first_run = None
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    if first_run is None:
                        first_run = run
                    else:
                        run.text = ""
            if first_run is not None:
                first_run.text = text
                return
        except Exception:
            pass
        try:
            shape.text = text
        except Exception:
            pass

    def remove_extra_slides(keep):
        sld_id_list = prs.slides._sldIdLst  # python-pptx has no public delete API.
        sld_ids = list(sld_id_list)
        for index in range(len(sld_ids) - 1, keep - 1, -1):
            rel_id = sld_ids[index].rId
            prs.part.drop_rel(rel_id)
            sld_id_list.remove(sld_ids[index])

    def title_for(index, slide_model):
        if index == 0:
            return trim(deck.get("title") or slide_model.get("title") or "Presentation", 26)
        return trim(slide_model.get("title") or f"Slide {index + 1}", 42)

    def body_for(slide_model):
        parts = []
        if slide_model.get("subtitle"):
            parts.append(slide_model.get("subtitle"))
        if slide_model.get("body"):
            parts.append(slide_model.get("body"))
        for bullet in slide_model.get("bullets") or []:
            parts.append(bullet)
        for column in slide_model.get("columns") or []:
            if column.get("title"):
                parts.append(column.get("title"))
            parts.extend(column.get("bullets") or [])
        if slide_model.get("kpi"):
            parts.append(f"{slide_model.get('kpi')} {slide_model.get('kpi_label', '')}".strip())
        if slide_model.get("quote"):
            parts.append(slide_model.get("quote"))
        return trim(" · ".join(parts), 260)

    def replace_cover(slide, slide_model):
        shapes = text_shapes(slide)
        if not shapes:
            return
        title = title_for(0, slide_model)
        subtitle = trim(deck.get("subtitle") or slide_model.get("subtitle") or "Q3 2026", 24)
        org = trim((deck.get("theme") or {}).get("organization") or deck.get("organization") or "", 24)
        # Largest textbox is usually the hero title.
        title_shape = max(shapes, key=lambda shape: int(shape.width) * int(shape.height))
        set_shape_text(title_shape, title)
        for shape in shapes:
            if shape is title_shape:
                continue
            txt = normalized(shape_text(shape))
            if re.search(r"\bq[1-4]\b|\b20\d{2}\b|date", txt):
                set_shape_text(shape, subtitle)
            elif org and len(txt) <= 40:
                set_shape_text(shape, org)
                org = ""

    def replace_agenda(slide, slides_models):
        shapes = text_shapes(slide)
        titles = [trim(item.get("title") or "", 30) for item in slides_models[1:5]]
        title_idx = 0
        for shape in shapes:
            txt = normalized(shape_text(shape))
            if txt in {"agenda", "table of contents", "overview"}:
                set_shape_text(shape, "Agenda")
                continue
            if re.fullmatch(r"\d{1,2}", txt):
                continue
            if title_idx < len(titles) and txt not in {"", "agenda"}:
                set_shape_text(shape, titles[title_idx])
                title_idx += 1

    def replace_content_slide(slide, index, slide_model):
        shapes = text_shapes(slide)
        if not shapes:
            return
        useful = [
            shape
            for shape in shapes
            if "back to agenda" not in normalized(shape_text(shape))
            and not re.fullmatch(r"\d{1,2}", normalized(shape_text(shape)))
        ]
        if not useful:
            return
        # Top-most sizeable text tends to be the slide heading.
        title_shape = min(
            useful,
            key=lambda shape: (int(shape.top), -(int(shape.width) * int(shape.height))),
        )
        set_shape_text(title_shape, title_for(index, slide_model))
        body = body_for(slide_model)
        if not body:
            return
        body_candidates = [shape for shape in useful if shape is not title_shape]
        if body_candidates:
            body_shape = max(body_candidates, key=lambda shape: int(shape.width) * int(shape.height))
            set_shape_text(body_shape, body)

    keep = min(len(slides_list), len(prs.slides))
    if keep < len(prs.slides):
        remove_extra_slides(keep)

    for index, slide_model in enumerate(slides_list[:keep]):
        slide = prs.slides[index]
        if index == 0:
            replace_cover(slide, slide_model)
        elif index == 1:
            replace_agenda(slide, slides_list)
        else:
            replace_content_slide(slide, index, slide_model)
        if slide_model.get("notes"):
            try:
                slide.notes_slide.notes_text_frame.text = str(slide_model.get("notes"))
            except Exception:
                pass

    prs.save(out_path)
    return {
        "logo": False,
        "img_ok": 0,
        "img_fail": 0,
        "notes": [f"template source preserved: {os.path.basename(source)}"],
    }


def render_pptx(deck, base_dir, out_path):
    """Build a native, editable .pptx from the deck model. Returns True on success."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except Exception as exc:  # python-pptx not installed → caller falls back to HTML/PDF
        sys.stderr.write(f"deck_render: python-pptx unavailable ({exc}); skipping .pptx\n")
        return False

    theme = {**DEFAULT_THEME, **(deck.get("theme") or {})}
    brand = RGBColor(*hexrgb(theme["primary"]))
    brand2 = RGBColor(*hexrgb(theme["secondary"], (26, 32, 44)))
    accent = RGBColor(*hexrgb(theme["accent"], (237, 137, 54)))
    ink = RGBColor(0x1A, 0x20, 0x2C)
    muted = RGBColor(0x4A, 0x55, 0x68)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    head_font = theme["heading_font"] or "Inter"
    body_font = theme["body_font"] or "Inter"

    # python-pptx needs real files; accept either a path or a data: URL (the brand-kit
    # logo arrives as a data URL). Materialised temp files are cleaned up after save.
    import tempfile

    _tmp_files = []

    def img_path(spec):
        if not spec:
            return ""
        if str(spec).startswith("data:"):
            try:
                header, b64 = spec.split(",", 1)
                ext = ".png"
                if "jpeg" in header or "jpg" in header:
                    ext = ".jpg"
                fd, path = tempfile.mkstemp(suffix=ext)
                with os.fdopen(fd, "wb") as fh:
                    fh.write(base64.b64decode(b64))
                _tmp_files.append(path)
                return path
            except Exception:
                return ""
        return resolve_path(spec, base_dir)

    logo_spec = theme.get("logo", "")
    logo_is_svg = isinstance(logo_spec, str) and (
        "image/svg" in logo_spec or logo_spec.lower().endswith(".svg")
    )
    logo_path = "" if logo_is_svg else img_path(logo_spec)
    # Observability: report exactly what got embedded so a missing logo/image is
    # never a silent mystery (the agent + the user see it in the run log).
    stats = {"logo": False, "img_ok": 0, "img_fail": 0, "notes": []}
    if logo_spec and logo_is_svg:
        stats["notes"].append(
            "logo is SVG — PowerPoint can't embed SVG; provide a PNG/JPG logo for the .pptx"
        )
    elif logo_spec and not logo_path:
        stats["notes"].append("logo could not be resolved (bad path or data URL)")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def fill_bg(slide, color):
        shp = slide.shapes.add_shape(1, 0, 0, SW, SH)  # rectangle
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        slide.shapes._spTree.remove(shp._element)
        slide.shapes._spTree.insert(2, shp._element)  # send to back
        return shp

    def accent_bar(slide):
        bar = slide.shapes.add_shape(1, 0, SH - Inches(0.14), SW, Inches(0.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False

    def add_logo(slide):
        if logo_path:
            try:
                slide.shapes.add_picture(
                    logo_path, SW - Inches(1.9), Inches(0.4), height=Inches(0.5)
                )
                stats["logo"] = True
            except Exception as exc:
                stats["notes"].append(f"logo add_picture failed: {exc}")

    def add_image(slide, spec, left, top, **kw):
        """add_picture with tracking; resolves a path or data URL via img_path."""
        path = img_path(spec)
        if not path:
            if spec:
                stats["img_fail"] += 1
                stats["notes"].append(f"image not resolved: {spec[:60]}")
            return
        try:
            slide.shapes.add_picture(path, left, top, **kw)
            stats["img_ok"] += 1
        except Exception as exc:
            stats["img_fail"] += 1
            stats["notes"].append(f"image add_picture failed ({spec[:40]}): {exc}")

    def accent_rail(slide):
        """Thin brand rail down the left edge — the signature of a content slide."""
        rail = slide.shapes.add_shape(1, 0, 0, Inches(0.16), SH)
        rail.fill.solid()
        rail.fill.fore_color.rgb = brand
        rail.line.fill.background()
        rail.shadow.inherit = False

    def title_underline(slide, left, top, width):
        rule = slide.shapes.add_shape(1, left, top, width, Pt(4))
        rule.fill.solid()
        rule.fill.fore_color.rgb = accent
        rule.line.fill.background()
        rule.shadow.inherit = False

    def footer(slide, page, total, org):
        label = f"{org}  ·  {page}/{total}" if org else f"{page}/{total}"
        textbox(slide, Inches(0.55), SH - Inches(0.55), Inches(6.0), Inches(0.4),
                [(label, 10, muted, body_font, False, False)])

    def textbox(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        first = True
        for (text, size, color, font, bold, bullet) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = ("•  " + text) if bullet else text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
            r.font.bold = bold
        return tb

    def notes(slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = str(text)

    org = deck.get("organization", "")
    slides_list = deck.get("slides", [])
    total = len(slides_list)
    for idx, s in enumerate(slides_list, 1):
        layout = s.get("layout", "bullets")
        slide = add_slide()
        title = s.get("title", "")

        if layout in ("cover", "section"):
            fill_bg(slide, brand)
            runs = [(title, 46 if layout == "cover" else 40, white, head_font, True, False)]
            if s.get("subtitle"):
                runs.append((s["subtitle"], 20, white, body_font, False, False))
            textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6), runs,
                    anchor=MSO_ANCHOR.MIDDLE)
            # accent rule under the title block
            rule = slide.shapes.add_shape(1, Inches(0.95), Inches(5.0), Inches(1.3), Pt(6))
            rule.fill.solid()
            rule.fill.fore_color.rgb = accent
            rule.line.fill.background()
            rule.shadow.inherit = False
            accent_bar(slide)
            add_logo(slide)
            notes(slide, s.get("notes"))
            continue

        # content slides: white bg, brand rail, underlined title, footer
        fill_bg(slide, white)
        accent_rail(slide)
        add_logo(slide)
        footer(slide, idx, total, org)
        if title:
            title_underline(slide, Inches(0.9), Inches(1.55), Inches(2.0))
            textbox(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.1),
                    [(title, 30, brand, head_font, True, False)])

        if layout == "kpi":
            textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2),
                    [(s.get("kpi", ""), 80, brand, head_font, True, False)],
                    anchor=MSO_ANCHOR.MIDDLE)
            textbox(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0),
                    [(s.get("kpi_label", ""), 22, muted, body_font, False, False)])
        elif layout == "quote":
            textbox(slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(2.6),
                    [("“" + s.get("quote", "") + "”", 32, brand, head_font, True, False)],
                    anchor=MSO_ANCHOR.MIDDLE)
            textbox(slide, Inches(1.2), Inches(4.9), Inches(10.9), Inches(0.8),
                    [("— " + s.get("author", ""), 18, muted, body_font, False, False)])
        elif layout in ("image_left", "image_right"):
            text_left = Inches(7.0) if layout == "image_left" else Inches(0.9)
            img_left = Inches(0.7) if layout == "image_left" else Inches(7.0)
            add_image(slide, s.get("image", ""), img_left, Inches(1.9), width=Inches(5.6))
            runs = [(b, 18, muted, body_font, False, True) for b in s.get("bullets", [])]
            if s.get("body"):
                runs.append((s["body"], 18, muted, body_font, False, False))
            if runs:
                textbox(slide, text_left, Inches(1.9), Inches(5.6), Inches(4.6), runs)
        elif layout == "two_column":
            cols = s.get("columns", [])[:2]
            for i, c in enumerate(cols):
                left = Inches(0.9) if i == 0 else Inches(7.0)
                runs = [(c.get("title", ""), 20, brand, head_font, True, False)]
                runs += [(b, 17, muted, body_font, False, True)
                         for b in c.get("bullets", [])]
                textbox(slide, left, Inches(1.9), Inches(5.6), Inches(4.6), runs)
        elif layout == "timeline":
            # Start/step/height tuned so 6 items (the max) fit inside the 7.5in slide:
            # 1.9 + 5*0.9 + 0.85 = 7.25in, leaving room for the footer/accent bar.
            top = 1.9
            for it in s.get("items", [])[:6]:
                textbox(slide, Inches(0.9), Inches(top), Inches(1.5), Inches(0.85),
                        [(it.get("label", ""), 18, brand, head_font, True, False)])
                dot = slide.shapes.add_shape(9, Inches(2.55), Inches(top + 0.08),
                                              Pt(11), Pt(11))  # 9 = MSO_SHAPE.OVAL
                dot.fill.solid()
                dot.fill.fore_color.rgb = accent
                dot.line.fill.background()
                dot.shadow.inherit = False
                runs = [(it.get("title", ""), 17, brand, head_font, True, False)]
                if it.get("text"):
                    runs.append((it.get("text", ""), 14, muted, body_font, False, False))
                textbox(slide, Inches(3.0), Inches(top), Inches(9.3), Inches(0.85), runs)
                top += 0.9
        elif layout == "comparison":
            headers = s.get("headers", [])[:4]
            rows = s.get("rows", [])[:8]
            if headers and rows:
                shape = slide.shapes.add_table(
                    len(rows) + 1, len(headers),
                    Inches(0.9), Inches(1.9), Inches(11.5),
                    Inches(min(4.6, 0.55 + 0.5 * len(rows))))
                table = shape.table
                for c, h in enumerate(headers):
                    table.cell(0, c).text = str(h)
                for r, row in enumerate(rows, start=1):
                    for c, cell_text in enumerate(row[: len(headers)]):
                        table.cell(r, c).text = str(cell_text)
                for row_cells in table.rows:
                    for cell in row_cells.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(13)
                                run.font.name = body_font
        elif layout == "team_grid":
            members = s.get("members", [])[:8]
            per_row = 4 if len(members) > 4 else max(len(members), 1)
            col_w = 11.5 / per_row
            for i, m in enumerate(members):
                row_i, col_i = divmod(i, per_row)
                left = 0.9 + col_i * col_w
                top = 2.1 + row_i * 2.3
                avatar = slide.shapes.add_shape(9, Inches(left), Inches(top),
                                                 Inches(0.85), Inches(0.85))
                avatar.fill.solid()
                avatar.fill.fore_color.rgb = brand
                avatar.line.fill.background()
                avatar.shadow.inherit = False
                avatar.text_frame.text = _initials(m.get("name", ""))
                textbox(slide, Inches(left), Inches(top + 0.95),
                        Inches(col_w - 0.3), Inches(1.1),
                        [(m.get("name", ""), 16, brand, head_font, True, False),
                         (m.get("role", ""), 13, muted, body_font, False, False)])
        else:  # bullets
            runs = [(b, 20, muted, body_font, False, True) for b in s.get("bullets", [])]
            if s.get("body"):
                runs.append((s["body"], 18, muted, body_font, False, False))
            has_img = bool(s.get("image"))
            text_w = Inches(7.4) if has_img else Inches(11.5)
            if runs:
                textbox(slide, Inches(0.9), Inches(1.9), text_w, Inches(4.6), runs)
            add_image(slide, s.get("image", ""), Inches(8.6), Inches(1.9), width=Inches(3.8))

        accent_bar(slide)
        notes(slide, s.get("notes"))

    prs.save(out_path)
    for tmp in _tmp_files:
        try:
            os.remove(tmp)
        except OSError:
            pass
    return stats


# ---------------------------------------------------------------------------- main
def main():
    ap = argparse.ArgumentParser(description="Render a Homun deck JSON to HTML + PPTX.")
    ap.add_argument("deck", nargs="?", help="path to deck.json")
    ap.add_argument("--prefix", default=None, help="output prefix (default: deck file stem)")
    ap.add_argument("--no-pptx", action="store_true", help="skip the .pptx output")
    ap.add_argument("--template-pptx", default=None, help="real .pptx/.potx template to edit for the .pptx output")
    ap.add_argument("--self-test", action="store_true", help="verify renderer quality contracts")
    args = ap.parse_args()

    if args.self_test:
        required = ["overflow-wrap:anywhere", "hyphens:auto"]
        missing = [item for item in required if item not in _HTML_CSS]
        print(json.dumps({"ok": not missing, "missing": missing}, ensure_ascii=False))
        return 0 if not missing else 2

    if not args.deck:
        ap.error("the following arguments are required: deck")

    with open(args.deck, "r", encoding="utf-8") as fh:
        deck = json.load(fh)
    if not deck.get("slides"):
        sys.exit("deck.json has no 'slides'.")

    base_dir = os.path.dirname(os.path.abspath(args.deck))

    # Auto-apply the brand from the output dir, if present. The gateway writes
    # `brand.json` (theme: colours/fonts/org) + `logo.png` next to deck.json when a deck
    # is generated, so the model can include ONLY slide content in deck.json — it never
    # has to embed the (large) logo data URL. A `theme` in deck.json still overrides.
    brand_file = os.path.join(base_dir, "brand.json")
    if os.path.isfile(brand_file):
        try:
            with open(brand_file, "r", encoding="utf-8") as fh:
                brand = json.load(fh)
            deck["theme"] = {**brand, **(deck.get("theme") or {})}
        except Exception:
            pass
    theme = deck.get("theme") or {}
    if not theme.get("logo") and os.path.isfile(os.path.join(base_dir, "logo.png")):
        theme["logo"] = "logo.png"
        deck["theme"] = theme

    prefix = args.prefix or os.path.splitext(os.path.basename(args.deck))[0]
    out_html = os.path.join(base_dir, f"{prefix}.html")
    out_pptx = os.path.join(base_dir, f"{prefix}.pptx")

    with open(out_html, "w", encoding="utf-8") as fh:
        fh.write(render_html(deck, base_dir))
    print(f"wrote {out_html}")

    if not args.no_pptx:
        result = (
            render_template_pptx(deck, base_dir, args.template_pptx, out_pptx)
            if args.template_pptx
            else render_pptx(deck, base_dir, out_pptx)
        )
        if result:
            print(f"wrote {out_pptx}")
            if args.template_pptx:
                print(f"  template: {args.template_pptx}")
            else:
                print(
                    f"  embedded: logo={'yes' if result['logo'] else 'NO'}, "
                    f"images={result['img_ok']} ok / {result['img_fail']} failed"
                )
            for note in result["notes"]:
                print(f"  ⚠ {note}")
        else:
            print("pptx skipped (python-pptx unavailable — restart the contained computer)")


if __name__ == "__main__":
    sys.exit(main())
